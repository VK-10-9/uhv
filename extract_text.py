import os
import sys
import glob
import re

# Try importing extraction libraries
try:
    from pypdf import PdfReader
except ImportError:
    print("Error: pypdf library is not installed. Please run: pip install pypdf")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("Warning: python-pptx library is not installed. PPTX extraction will be skipped. Run: pip install python-pptx")
    Presentation = None

# For older .ppt files, try win32com on Windows
pywin32_available = False
try:
    import win32com.client
    pywin32_available = True
except ImportError:
    print("Warning: pywin32 is not installed. Older .ppt extraction will be skipped. Run: pip install pywin32")

def extract_pdf(filepath):
    print(f"Extracting PDF: {os.path.basename(filepath)}...")
    text_content = []
    try:
        reader = PdfReader(filepath)
        for idx, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                text_content.append(f"--- Page {idx + 1} ---\n{text.strip()}")
    except Exception as e:
        print(f"Failed to extract {filepath}: {e}")
    return "\n\n".join(text_content)

def extract_pptx(filepath):
    if Presentation is None:
        return ""
    print(f"Extracting PPTX: {os.path.basename(filepath)}...")
    text_content = []
    try:
        prs = Presentation(filepath)
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_content.append(f"--- Slide {idx + 1} ---\n" + "\n".join(slide_text))
    except Exception as e:
        print(f"Failed to extract {filepath}: {e}")
    return "\n\n".join(text_content)

def extract_ppt(filepath, powerpoint_app=None):
    if not pywin32_available or powerpoint_app is None:
        return ""
    print(f"Extracting PPT via PowerPoint COM: {os.path.basename(filepath)}...")
    text_content = []
    try:
        abs_path = os.path.abspath(filepath)
        # Open PowerPoint presentation in read-only and without window
        # 1 = ReadOnly, 0 = Untitled (False), 0 = WithWindow (False)
        pres = powerpoint_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)
        for idx, slide in enumerate(pres.Slides):
            slide_text = []
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    text = shape.TextFrame.TextRange.Text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                text_content.append(f"--- Slide {idx + 1} ---\n" + "\n".join(slide_text))
        pres.Close()
    except Exception as e:
        print(f"Failed to extract {filepath} via COM: {e}")
    return "\n\n".join(text_content)

def main():
    workspace = "."
    output_file = "extracted_notes.txt"
    
    # Initialize PowerPoint COM app if pywin32 is available
    ppt_app = None
    if pywin32_available:
        try:
            # Try to get or launch PowerPoint application
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.DisplayAlerts = False
        except Exception as e:
            print(f"Warning: Could not launch PowerPoint application via COM: {e}. .ppt files will be skipped.")
            ppt_app = None

    all_files = glob.glob(os.path.join(workspace, "*"))
    extracted_data = []

    for filepath in all_files:
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()

        if filename == output_file or filename == "extract_text.py":
            continue

        content = ""
        if ext == ".pdf":
            content = extract_pdf(filepath)
        elif ext == ".pptx":
            content = extract_pptx(filepath)
        elif ext == ".ppt":
            if ppt_app:
                content = extract_ppt(filepath, ppt_app)
            else:
                print(f"Skipping older PPT (no PowerPoint COM): {filename}")
        else:
            continue

        if content:
            header = f"=========================================\nSOURCE FILE: {filename}\n========================================="
            extracted_data.append(f"{header}\n\n{content}\n")

    # Close PowerPoint app if opened
    if ppt_app:
        try:
            ppt_app.Quit()
        except:
            pass

    if extracted_data:
        with open(output_file, "w", encoding="utf-8") as f:
            f.write("\n\n".join(extracted_data))
        print(f"\nSuccess! Extracted text written to '{output_file}'")
    else:
        print("No readable PDF, PPT, or PPTX files found in this directory.")

if __name__ == "__main__":
    main()
